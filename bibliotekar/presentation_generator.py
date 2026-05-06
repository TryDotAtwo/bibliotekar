"""
entity_id=presentation_generator_module; type=module; state=initial

Module for generating slide decks from wiki content or query results. Uses python‑pptx to produce a PPTX file. The function ``create_presentation`` accepts a title and a list of sections (heading + body) and writes a file to disk. The caller is responsible for syncing the generated file with the user.
"""

from pathlib import Path
from typing import Iterable, Tuple

from pptx import Presentation  # type: ignore
from pptx.util import Inches


def create_presentation(title: str, sections: Iterable[Tuple[str, str]], output_path: Path) -> Path:
    """Generate a PowerPoint presentation.

    entity_id=create_presentation_function; type=function; state=initial
    condition=sections_provided → action=generate_slides → result=pptx_file

    Parameters
    ----------
    title: str
        Title slide heading.
    sections: Iterable[Tuple[str, str]]
        Sequence of (heading, body) for each slide.
    output_path: Path
        Destination path for the generated PPTX.

    Returns
    -------
    Path
        Path to the generated PPTX file.
    """
    pres = Presentation()
    # Title slide
    title_slide_layout = pres.slide_layouts[0]
    slide = pres.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated on demand"
    # Content slides
    bullet_layout = pres.slide_layouts[1]
    for heading, body in sections:
        slide = pres.slides.add_slide(bullet_layout)
        slide.shapes.title.text = heading
        tf = slide.shapes.placeholders[1].text_frame
        for line in body.split("\n"):
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    pres.save(str(output_path))
    return output_path